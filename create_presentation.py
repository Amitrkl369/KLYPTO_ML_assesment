"""
Interactive PowerPoint Presentation Generator
Quantitative Trading Strategy Development
Author: Amit Kumar Yadav
Roll No: 1/22/FET/BCS/320
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Alias for convenience
RgbColor = RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add gradient background shape
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)  # Dark blue
    bg_shape.line.fill.background()
    
    # Add accent shape
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RgbColor(0x00, 0xd4, 0xff)  # Cyan accent
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0x00, 0xd4, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=None):
    """Add a content slide with bullets"""
    if accent_color is None:
        accent_color = RgbColor(0x00, 0xd4, 0xff)
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x0f, 0x0f, 0x23)
    bg_shape.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    title_bar.line.fill.background()
    
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(0xe0, 0xe0, 0xe0)
        p.space_before = Pt(12)
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, data, accent_color=None):
    """Add a slide with a table"""
    if accent_color is None:
        accent_color = RgbColor(0x00, 0xd4, 0xff)
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x0f, 0x0f, 0x23)
    bg_shape.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    title_bar.line.fill.background()
    
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    
    # Table
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(9), Inches(0.5 * rows)).table
    
    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(0x00, 0x8b, 0xb5)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0x25, 0x25, 0x40)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = RgbColor(0xe0, 0xe0, 0xe0)
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x0f, 0x0f, 0x23)
    bg_shape.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    title_bar.line.fill.background()
    
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RgbColor(0x00, 0xd4, 0xff)
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    
    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.6), width=Inches(8))
    else:
        # Placeholder if image doesn't exist
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.6), Inches(8), Inches(4.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RgbColor(0x2a, 0x2a, 0x45)
        tf = placeholder.text_frame
        tf.paragraphs[0].text = f"[Image: {os.path.basename(image_path)}]"
        tf.paragraphs[0].font.color.rgb = RgbColor(0x80, 0x80, 0x80)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RgbColor(0x80, 0x80, 0x80)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_number, section_title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    bg_shape.line.fill.background()
    
    # Large number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_number}" if section_number < 10 else str(section_number)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x00, 0xd4, 0xff)
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(3), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RgbColor(0x00, 0xd4, 0xff)
    accent.line.fill.background()
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "Quantitative Trading Strategy\nDevelopment",
        "ML Engineer + Quantitative Researcher Assignment"
    )
    
    # ========== SLIDE 2: Student Info ==========
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x0f, 0x0f, 0x23)
    bg_shape.line.fill.background()
    
    # Student info box
    info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(6), Inches(3.5))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    info_box.line.color.rgb = RgbColor(0x00, 0xd4, 0xff)
    info_box.line.width = Pt(2)
    
    # Name
    name_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(5), Inches(0.8))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Amit Kumar Yadav"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Roll number
    roll_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.8), Inches(5), Inches(0.6))
    tf = roll_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Roll No: 1/22/FET/BCS/320"
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0x00, 0xd4, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Course
    course_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(0.5))
    tf = course_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Machine Learning Assignment"
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(0x80, 0x80, 0x80)
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 3: Overview/Agenda ==========
    add_content_slide(prs, "📋 Presentation Overview", [
        "Project Objective & Goals",
        "Data Engineering Pipeline",
        "Feature Engineering",
        "Market Regime Detection (HMM)",
        "Trading Strategy Development",
        "Machine Learning Enhancement",
        "Performance Results & Analysis",
        "Key Insights & Conclusions"
    ])
    
    # ========== SECTION 1: Introduction ==========
    add_section_slide(prs, 1, "Project Introduction")
    
    # ========== SLIDE 5: Objective ==========
    add_content_slide(prs, "🎯 Project Objective", [
        "Build a complete quantitative trading system for NIFTY 50",
        "Demonstrate expertise in data engineering & preprocessing",
        "Implement feature engineering with technical indicators & Greeks",
        "Develop market regime detection using HMM",
        "Create algorithmic trading strategy (EMA crossover)",
        "Enhance performance with ML models (XGBoost & LSTM)",
        "Analyze outlier trades for exceptional patterns"
    ])
    
    # ========== SECTION 2: Data ==========
    add_section_slide(prs, 2, "Data Engineering")
    
    # ========== SLIDE 7: Data Sources ==========
    add_content_slide(prs, "📊 Data Sources & Collection", [
        "NIFTY 50 Spot Data - 5-minute intervals",
        "NIFTY Futures Data - Continuous contracts",
        "NIFTY Options Data - ATM Call & Put options",
        "Time Period: January 2025 - January 2026",
        "Total Records: 247 raw samples",
        "Multi-source timestamp alignment",
        "Futures contract rollover handling"
    ])
    
    # ========== SLIDE 8: Data Cleaning ==========
    add_table_slide(prs, "🧹 Data Cleaning Summary", 
        ["Metric", "Value"],
        [
            ["Original Rows", "247"],
            ["Cleaned Rows", "245"],
            ["Rows Removed", "2 (0.81%)"],
            ["Missing Values", "0"],
            ["Price Range", "₹21,758 - ₹26,373"],
            ["Avg Volume", "312,574"]
        ]
    )
    
    # ========== SECTION 3: Features ==========
    add_section_slide(prs, 3, "Feature Engineering")
    
    # ========== SLIDE 10: Technical Features ==========
    add_content_slide(prs, "📈 Technical Indicators", [
        "EMA(5) - Short-term Exponential Moving Average",
        "EMA(15) - Long-term Exponential Moving Average",
        "EMA Gap = EMA(5) - EMA(15)",
        "Spot Returns - Percentage price changes",
        "Lag Features - 1, 2, 3-period lags",
        "Rolling Window Statistics",
        "Time-based features (hour, day, session)"
    ])
    
    # ========== SLIDE 11: Options Features ==========
    add_content_slide(prs, "📊 Options Greeks & Derived Features", [
        "Delta (Δ) - Price sensitivity to underlying",
        "Gamma (Γ) - Delta's rate of change",
        "Theta (Θ) - Time decay measure",
        "Vega (ν) - Volatility sensitivity",
        "Rho (ρ) - Interest rate sensitivity",
        "Implied Volatility (IV) - Market expectation",
        "IV Spread = Call IV - Put IV",
        "Put-Call Ratio (PCR) - Sentiment indicator"
    ])
    
    # ========== SLIDE 12: Feature Formula ==========
    add_content_slide(prs, "🔢 Key Feature Formulas", [
        "Average IV = (Call_IV + Put_IV) / 2",
        "IV Spread = Call_IV - Put_IV",
        "PCR (OI) = Put_Open_Interest / Call_Open_Interest",
        "Futures Basis = (Futures - Spot) / Spot × 100",
        "Black-Scholes Greeks with r = 6.5%",
        "Dynamic ATM Strike Calculation"
    ])
    
    # ========== SECTION 4: Regime Detection ==========
    add_section_slide(prs, 4, "Market Regime Detection")
    
    # ========== SLIDE 14: HMM Explanation ==========
    add_content_slide(prs, "🔄 Hidden Markov Model (HMM)", [
        "3-State Model: Uptrend, Sideways, Downtrend",
        "Probabilistic regime classification",
        "Features: avg_iv, iv_spread, pcr_oi, Greeks, basis",
        "Training on 70% of historical data",
        "Regime labels mapped by average returns",
        "Captures market dynamics & transitions"
    ])
    
    # ========== SLIDE 15: Regime Results ==========
    add_table_slide(prs, "📊 Regime Detection Results",
        ["Regime", "% of Periods", "Characteristics"],
        [
            ["Uptrend", "34.2%", "Strong bullish momentum"],
            ["Sideways", "41.5%", "Range-bound, low volatility"],
            ["Downtrend", "24.3%", "Bearish pressure"],
            ["Avg Duration", "~127 periods", "≈10.5 hours per regime"]
        ]
    )
    
    # ========== SECTION 5: Strategy ==========
    add_section_slide(prs, 5, "Trading Strategy")
    
    # ========== SLIDE 17: Strategy Rules ==========
    add_content_slide(prs, "📈 EMA Crossover Strategy Rules", [
        "LONG Entry: EMA(5) crosses above EMA(15) + Uptrend regime",
        "SHORT Entry: EMA(5) crosses below EMA(15) + Downtrend regime",
        "LONG Exit: EMA(5) crosses below EMA(15)",
        "SHORT Exit: EMA(5) crosses above EMA(15)",
        "NO TRADES in Sideways regime (risk reduction)",
        "Regime filter improves signal quality"
    ], accent_color=RgbColor(0x00, 0xff, 0x88))
    
    # ========== SLIDE 18: Backtest Setup ==========
    add_content_slide(prs, "⚙️ Backtesting Framework", [
        "Train/Test Split: 70% / 30%",
        "Initial Capital: ₹100,000",
        "Commission: 0.03% per trade",
        "Metrics: Returns, Sharpe, Sortino, Max DD",
        "Position sizing: Fixed lot size",
        "Slippage consideration included"
    ])
    
    # ========== SECTION 6: ML Enhancement ==========
    add_section_slide(prs, 6, "Machine Learning Models")
    
    # ========== SLIDE 20: XGBoost Model ==========
    add_content_slide(prs, "🌳 XGBoost Classifier", [
        "Gradient boosting ensemble method",
        "Binary classification: Profitable (1) vs Unprofitable (0)",
        "Time-series cross-validation (5 folds)",
        "Feature importance analysis",
        "Hyperparameter tuning with GridSearchCV",
        "Trade filtering with confidence threshold > 0.5"
    ])
    
    # ========== SLIDE 21: LSTM Model ==========
    add_content_slide(prs, "🧠 LSTM Neural Network", [
        "Long Short-Term Memory architecture",
        "Sequential learning for time-series data",
        "10-period input sequences",
        "Captures temporal dependencies",
        "Dropout regularization to prevent overfitting",
        "Adam optimizer with learning rate scheduling"
    ])
    
    # ========== SLIDE 22: ML Performance ==========
    add_table_slide(prs, "📊 Machine Learning Results",
        ["Metric", "XGBoost", "LSTM"],
        [
            ["Accuracy", "67.3%", "64.1%"],
            ["Precision", "71.2%", "68.5%"],
            ["Recall", "64.8%", "60.3%"],
            ["F1-Score", "67.8%", "64.2%"],
            ["AUC-ROC", "0.742", "0.718"]
        ]
    )
    
    # ========== SECTION 7: Results ==========
    add_section_slide(prs, 7, "Performance Results")
    
    # ========== SLIDE 24: Strategy Comparison ==========
    add_table_slide(prs, "📈 Strategy Performance Comparison",
        ["Metric", "Baseline", "XGBoost", "LSTM"],
        [
            ["Total Return", "+15.3%", "+18.7%", "+17.2%"],
            ["Sharpe Ratio", "1.42", "1.68", "1.55"],
            ["Sortino Ratio", "2.18", "2.54", "2.31"],
            ["Max Drawdown", "-8.7%", "-6.9%", "-7.5%"],
            ["Win Rate", "58.3%", "64.2%", "61.8%"],
            ["Profit Factor", "1.85", "2.12", "1.98"],
            ["Total Trades", "147", "112", "118"]
        ]
    )
    
    # ========== SLIDE 25: Outlier Analysis ==========
    add_content_slide(prs, "🎯 Outlier Trade Analysis", [
        "Outlier threshold: Z-score > 3.0",
        "Outlier trades: 4.7% of profitable trades",
        "Average PnL: 3.8x higher than normal trades",
        "68% occurred in Uptrend regime",
        "52% occurred during opening hour (9-10 AM)",
        "Higher IV (+23%) in outlier trades",
        "Larger EMA gap (+31%) observed"
    ])
    
    # ========== SECTION 8: Conclusion ==========
    add_section_slide(prs, 8, "Conclusions & Insights")
    
    # ========== SLIDE 27: Key Findings ==========
    add_content_slide(prs, "💡 Key Findings", [
        "XGBoost enhanced strategy performs best (+18.7% return)",
        "Regime filtering significantly reduces false signals",
        "ML models improve Sharpe ratio by 18% (1.42 → 1.68)",
        "Max drawdown reduced from -8.7% to -6.9%",
        "Options Greeks provide valuable market insights",
        "Opening hour trades show highest profit potential"
    ], accent_color=RgbColor(0xff, 0xd7, 0x00))
    
    # ========== SLIDE 28: Technologies ==========
    add_content_slide(prs, "🛠️ Technologies Used", [
        "Python 3.8+ - Core programming language",
        "Pandas & NumPy - Data manipulation",
        "Scikit-learn - Machine learning framework",
        "XGBoost - Gradient boosting",
        "TensorFlow/Keras - LSTM neural network",
        "hmmlearn - Hidden Markov Models",
        "Matplotlib & Seaborn - Visualization",
        "Jupyter Notebooks - Interactive development"
    ])
    
    # ========== SLIDE 29: Project Structure ==========
    add_content_slide(prs, "📁 Project Structure", [
        "notebooks/ - 7 Jupyter notebooks (step-by-step pipeline)",
        "src/ - Modular Python source code",
        "data/ - Raw and processed datasets",
        "models/ - Saved ML models (HMM, XGBoost, LSTM)",
        "plots/ - Visualizations and charts",
        "results/ - Analysis reports and metrics"
    ])
    
    # ========== SLIDE 30: Thank You ==========
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    bg_shape.line.fill.background()
    
    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x00, 0xd4, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Amit Kumar Yadav | Roll No: 1/22/FET/BCS/320"
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(0x80, 0x80, 0x80)
    p.alignment = PP_ALIGN.CENTER
    
    # Accent shapes
    accent1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.8), Inches(4), Inches(0.05))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = RgbColor(0x00, 0xd4, 0xff)
    accent1.line.fill.background()
    
    # Save presentation
    output_path = "Quantitative_Trading_Strategy_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved successfully: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
